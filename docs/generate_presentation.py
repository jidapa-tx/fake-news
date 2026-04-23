"""
Generate ชัวร์ก่อนแชร์ presentation as PowerPoint (.pptx)
Run: python3 docs/generate_presentation.py
Output: docs/presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Color palette ────────────────────────────────────────────────────────────
BG       = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
SURFACE  = RGBColor(0x1A, 0x2D, 0x44)   # dark blue-grey
ACCENT   = RGBColor(0x4F, 0x9C, 0xF9)   # sky blue
ACCENT2  = RGBColor(0x6E, 0xC6, 0xFF)   # light blue
TEXT     = RGBColor(0xE8, 0xF1, 0xFF)   # text primary
MUTED    = RGBColor(0x7A, 0x9C, 0xC4)   # text muted
RED      = RGBColor(0xEF, 0x44, 0x44)
ORANGE   = RGBColor(0xF9, 0x73, 0x16)
YELLOW   = RGBColor(0xEA, 0xB3, 0x08)
LIME     = RGBColor(0x84, 0xCC, 0x16)
GREEN    = RGBColor(0x22, 0xC5, 0x5E)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ──────────────────────────────────────────────────────────────────

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]  # completely blank
    slide = prs.slides.add_slide(blank_layout)
    # Fill background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide


def add_textbox(slide, left, top, width, height, text,
                font_size=18, bold=False, color=TEXT,
                align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txBox


def add_title(slide, title: str, subtitle: str = ""):
    # Accent bar at top
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), SLIDE_W, Pt(6)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    add_textbox(slide,
                left=Inches(0.5), top=Inches(0.15),
                width=Inches(12.3), height=Inches(0.7),
                text=title, font_size=30, bold=True,
                color=TEXT, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide,
                    left=Inches(0.5), top=Inches(0.85),
                    width=Inches(12.3), height=Inches(0.4),
                    text=subtitle, font_size=16, bold=False,
                    color=MUTED, align=PP_ALIGN.LEFT)


def add_card(slide, left, top, width, height, color=SURFACE):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = ACCENT
    shape.line.width = Pt(0.75)
    return shape


def add_bullets(slide, left, top, width, height, items: list,
                font_size=17, color=TEXT, muted_color=MUTED,
                indent_char="• "):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        if isinstance(item, tuple):
            run.text = indent_char + item[0]
            run.font.color.rgb = item[1]
        else:
            run.text = indent_char + item
            run.font.color.rgb = color
        run.font.size = Pt(font_size)
    return txBox


def add_notes(slide, text: str):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


# ── Slide builders ───────────────────────────────────────────────────────────

def slide_01_cover(prs):
    slide = blank_slide(prs)

    # Large accent circle decoration
    circle = slide.shapes.add_shape(9, Inches(9.5), Inches(-1), Inches(5), Inches(5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT
    circle.line.fill.background()
    # Make it semi-transparent via XML hack
    spPr = circle._element.spPr
    from lxml import etree
    solidFill = spPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    if solidFill is not None:
        srgbClr = solidFill.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        if srgbClr is not None:
            alpha = etree.SubElement(srgbClr, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            alpha.set('val', '10000')  # 10% opacity

    add_textbox(slide,
                left=Inches(0.7), top=Inches(1.5),
                width=Inches(9), height=Inches(1.2),
                text="ชัวร์ก่อนแชร์", font_size=56, bold=True,
                color=ACCENT, align=PP_ALIGN.LEFT)

    add_textbox(slide,
                left=Inches(0.7), top=Inches(2.9),
                width=Inches(9), height=Inches(0.7),
                text="Sure Before Share — Thai Fake News Detector", font_size=22,
                color=ACCENT2, align=PP_ALIGN.LEFT)

    add_textbox(slide,
                left=Inches(0.7), top=Inches(3.7),
                width=Inches(9), height=Inches(0.6),
                text="ตรวจสอบข่าวก่อนส่งต่อ — เพราะข่าวปลอมแพร่เร็วกว่าที่คุณคิด",
                font_size=18, color=MUTED, align=PP_ALIGN.LEFT)

    # Divider
    line = slide.shapes.add_shape(1, Inches(0.7), Inches(4.45), Inches(4), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    add_textbox(slide,
                left=Inches(0.7), top=Inches(4.6),
                width=Inches(6), height=Inches(0.5),
                text="ทีม: [ชื่อทีม]   |   วันที่: เมษายน 2569", font_size=15,
                color=MUTED, align=PP_ALIGN.LEFT)

    add_textbox(slide,
                left=Inches(0.7), top=Inches(5.2),
                width=Inches(6), height=Inches(0.5),
                text="Stack: Next.js 14 · TypeScript · PostgreSQL · Google Gemini API",
                font_size=13, color=MUTED, align=PP_ALIGN.LEFT)

    add_notes(slide,
              "ยินดีต้อนรับทุกท่านสู่การนำเสนอโปรเจกต์ ชัวร์ก่อนแชร์ "
              "แอปพลิเคชันตรวจสอบข่าวปลอมสำหรับผู้ใช้ภาษาไทย "
              "วันนี้เราจะพาทุกท่านไปรู้จักกับปัญหา วิธีแก้ไข และสถาปัตยกรรมทางเทคนิคของโปรเจกต์นี้")


def slide_02_problem(prs):
    slide = blank_slide(prs)
    add_title(slide, "ปัญหา: ข่าวปลอมในยุคดิจิทัล",
              "Problem Statement")

    add_bullets(slide,
                left=Inches(0.6), top=Inches(1.4),
                width=Inches(7.5), height=Inches(5),
                items=[
                    "ข่าวปลอมแพร่กระจายเร็วกว่าข่าวจริง 6 เท่า (MIT, 2018)",
                    "LINE/Facebook เป็นช่องทางหลักในการแชร์ข่าวปลอมในไทย",
                    "ผู้สูงอายุและผู้ใช้ทั่วไปขาดทักษะตรวจสอบข้อมูล",
                    "เครื่องมือตรวจสอบที่มีอยู่ต้องการการสมัครสมาชิก",
                    "ไม่มีเครื่องมือที่ใส่ใจความเป็นส่วนตัวของผู้ใช้",
                    "ขาดการรองรับภาษาไทยในเครื่องมือระดับสากล",
                ],
                font_size=19)

    # Right: stat cards
    for i, (stat, label, clr) in enumerate([
        ("6×", "เร็วกว่าข่าวจริง", RED),
        ("82%", "ผู้ใช้ไทยพบข่าวปลอม\nบน social media", ORANGE),
        ("0", "เครื่องมือฟรี ไม่ต้องสมัคร\nที่รองรับภาษาไทย", ACCENT),
    ]):
        card_top = Inches(1.5 + i * 1.85)
        add_card(slide, Inches(8.4), card_top, Inches(4.3), Inches(1.6))
        add_textbox(slide, Inches(8.5), card_top + Pt(8), Inches(4.1), Inches(0.7),
                    stat, font_size=36, bold=True, color=clr, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(8.5), card_top + Inches(0.85), Inches(4.1), Inches(0.6),
                    label, font_size=14, color=MUTED, align=PP_ALIGN.CENTER)

    add_notes(slide,
              "ข่าวปลอมเป็นปัญหาระดับชาติที่ส่งผลกระทบต่อสังคมไทยอย่างมาก "
              "โดยเฉพาะในช่วงวิกฤต เช่น COVID-19 หรือช่วงการเลือกตั้ง "
              "ปัจจุบันยังไม่มีเครื่องมือที่ใช้งานง่าย ฟรี และรักษาความเป็นส่วนตัวของผู้ใช้")


def slide_03_solution(prs):
    slide = blank_slide(prs)
    add_title(slide, "วิธีแก้: ชัวร์ก่อนแชร์", "Solution Overview")

    add_textbox(slide,
                left=Inches(0.6), top=Inches(1.35),
                width=Inches(12), height=Inches(1.0),
                text="แอปพลิเคชันเว็บสำหรับตรวจสอบความน่าเชื่อถือของข้อมูล "
                     "รองรับทั้งข้อความ URL และรูปภาพ โดยไม่ต้องสมัครสมาชิก "
                     "และไม่เก็บข้อมูลส่วนตัวของผู้ใช้",
                font_size=19, color=TEXT)

    promises = [
        ("🚫 ไม่ต้อง Login", "เข้าใช้งานได้ทันที ไม่มีกำแพงสมัครสมาชิก"),
        ("🔒 ไม่เก็บข้อมูล", "ประวัติการค้นหาเก็บในเครื่องของคุณเท่านั้น"),
        ("⚡ ตรวจสอบใน วินาที", "วิเคราะห์ด้วย AI + ฐานข้อมูลแหล่งข่าวน่าเชื่อถือ"),
        ("🌐 รองรับภาษาไทย", "ผลลัพธ์และเหตุผลเป็นภาษาไทย พร้อม EN fallback"),
    ]

    for i, (title_text, desc) in enumerate(promises):
        col = i % 2
        row = i // 2
        cx = Inches(0.6 + col * 6.3)
        cy = Inches(2.65 + row * 2.0)
        add_card(slide, cx, cy, Inches(5.9), Inches(1.7))
        add_textbox(slide, cx + Inches(0.2), cy + Inches(0.15),
                    Inches(5.5), Inches(0.55),
                    title_text, font_size=20, bold=True, color=ACCENT)
        add_textbox(slide, cx + Inches(0.2), cy + Inches(0.75),
                    Inches(5.5), Inches(0.7),
                    desc, font_size=16, color=MUTED)

    add_notes(slide,
              "ชัวร์ก่อนแชร์ถูกออกแบบมาสำหรับผู้ใช้ทั่วไปที่ต้องการตรวจสอบข้อมูลอย่างรวดเร็ว "
              "หลักการสำคัญคือ Privacy First — ไม่มีข้อมูลส่วนตัวออกจากเครื่องของผู้ใช้ "
              "และการวิเคราะห์ให้ผลลัพธ์ที่เข้าใจง่ายในภาษาไทย")


def slide_04_features(prs):
    slide = blank_slide(prs)
    add_title(slide, "ฟีเจอร์หลัก (FR.1–FR.7)", "Key Features")

    features = [
        ("FR.1", "🔓", "Privacy First", "ไม่ต้อง Login · ประวัติใน localStorage เท่านั้น"),
        ("FR.2", "🔍", "วิเคราะห์หลายรูปแบบ", "ข้อความ / URL / รูปภาพ · คะแนน 0–100 · 5 ระดับ"),
        ("FR.3", "📚", "รายการอ้างอิง", "แหล่งข่าวน่าเชื่อถือ 20+ แหล่ง · badge stance"),
        ("FR.4", "🌐", "วิเคราะห์โดเมน", "อายุโดเมน · SSL · WHOIS · phishing detection"),
        ("FR.5", "🖼️", "ตรวจรูปภาพ AI", "Gemini detect AI-generated · EXIF metadata"),
        ("FR.6", "🔥", "Trending", "ค้นหายอดนิยม · anonymized · SHA-256 hash"),
        ("FR.7", "📋", "ประวัติการค้นหา", "บันทึกใน browser · ส่งออก/นำเข้า JSON"),
    ]

    for i, (fr, icon, name, desc) in enumerate(features):
        col = i % 2
        row = i // 2
        if i == 6:  # last item centered
            cx = Inches(3.5)
        else:
            cx = Inches(0.4 + col * 6.5)
        cy = Inches(1.4 + row * 1.5)
        add_card(slide, cx, cy, Inches(6.0), Inches(1.3))
        add_textbox(slide, cx + Inches(0.15), cy + Inches(0.08),
                    Inches(0.6), Inches(0.5),
                    icon, font_size=22, color=TEXT)
        add_textbox(slide, cx + Inches(0.75), cy + Inches(0.08),
                    Inches(1.0), Inches(0.4),
                    fr, font_size=12, bold=True, color=ACCENT)
        add_textbox(slide, cx + Inches(0.75), cy + Inches(0.45),
                    Inches(5.0), Inches(0.35),
                    name, font_size=17, bold=True, color=TEXT)
        add_textbox(slide, cx + Inches(0.75), cy + Inches(0.82),
                    Inches(5.0), Inches(0.35),
                    desc, font_size=13, color=MUTED)

    add_notes(slide,
              "ฟีเจอร์ทั้ง 7 ถูกออกแบบให้ครอบคลุมการตรวจสอบทุกรูปแบบ "
              "ตั้งแต่ข้อความง่ายๆ ไปจนถึงรูปภาพที่อาจถูกสร้างด้วย AI "
              "โดยให้ความสำคัญกับความเป็นส่วนตัวและความง่ายในการใช้งาน")


def slide_05_demo_home(prs):
    slide = blank_slide(prs)
    add_title(slide, "Demo: หน้าหลัก (Home)", "Screen walkthrough")

    add_bullets(slide,
                left=Inches(0.5), top=Inches(1.4),
                width=Inches(7.5), height=Inches(5.5),
                items=[
                    "🏠 Hero Banner — ชื่อแอป + tagline + ไอคอนโล่",
                    "📝 Input Card — 3 แท็บ: ข้อความ / URL / รูปภาพ",
                    "📊 Stats Row — จำนวนการตรวจสอบวันนี้ + ความแม่นยำ",
                    "🕐 Recent Checks — รายการล่าสุดจาก localStorage",
                    "🔥 Trending Bar — ลิงก์ไปหน้า Trending",
                    "🌐 Language Toggle — TH / EN switcher",
                    "ปุ่ม Submit ทริกเกอร์ POST /api/analyze/text หรือ /image",
                ],
                font_size=17)

    # Phone frame (CSS-like placeholder using shapes)
    phone_left = Inches(8.3)
    phone_top  = Inches(1.1)
    phone_w    = Inches(3.5)
    phone_h    = Inches(5.9)

    frame = slide.shapes.add_shape(1, phone_left, phone_top, phone_w, phone_h)
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)
    frame.line.color.rgb = ACCENT
    frame.line.width = Pt(2.5)

    # Notch
    notch = slide.shapes.add_shape(9, phone_left + Inches(1.2), phone_top + Inches(0.1),
                                   Inches(1.1), Inches(0.25))
    notch.fill.solid()
    notch.fill.fore_color.rgb = ACCENT
    notch.line.fill.background()

    add_textbox(slide,
                left=phone_left + Inches(0.1), top=phone_top + Inches(2.2),
                width=phone_w - Inches(0.2), height=Inches(1.4),
                text="📱\nวางภาพ screenshot\nที่นี่",
                font_size=14, color=MUTED, align=PP_ALIGN.CENTER)

    add_notes(slide,
              "หน้าหลักถูกออกแบบให้เรียบง่ายและใช้งานได้ทันที "
              "ผู้ใช้สามารถวางข้อความ URL หรืออัปโหลดรูปภาพเพื่อตรวจสอบได้เลย "
              "โดยไม่ต้องสมัครสมาชิกหรือรอโหลดนาน")


def slide_06_demo_result(prs):
    slide = blank_slide(prs)
    add_title(slide, "Demo: หน้าผลลัพธ์ (Result)", "Screen walkthrough")

    add_bullets(slide,
                left=Inches(0.5), top=Inches(1.4),
                width=Inches(7.5), height=Inches(5.5),
                items=[
                    "🎯 Verdict Card — วงกลมคะแนน + 5-segment bar",
                    "   ระดับ: อันตราย / น่าสงสัย / ไม่แน่ใจ / ค่อนข้างจริง / ยืนยันแล้ว",
                    "📊 Evidence Stats — 3 การ์ด: แหล่งสนับสนุน / ต่อต้าน / AI confidence",
                    "📖 Reasoning Accordion — เหตุผล AI เป็น bullet ภาษาไทย",
                    "📰 Reference List — badge stance (สนับสนุน/ต่อต้าน/เป็นกลาง)",
                    "🌐 Source Analysis — ความน่าเชื่อถือโดเมน + SSL + WHOIS",
                    "📤 Action Bar — Copy Link / Share / ตรวจสอบใหม่",
                ],
                font_size=16)

    phone_left = Inches(8.3)
    phone_top  = Inches(1.1)
    phone_w    = Inches(3.5)
    phone_h    = Inches(5.9)

    frame = slide.shapes.add_shape(1, phone_left, phone_top, phone_w, phone_h)
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)
    frame.line.color.rgb = ACCENT
    frame.line.width = Pt(2.5)

    notch = slide.shapes.add_shape(9, phone_left + Inches(1.2), phone_top + Inches(0.1),
                                   Inches(1.1), Inches(0.25))
    notch.fill.solid()
    notch.fill.fore_color.rgb = ACCENT
    notch.line.fill.background()

    add_textbox(slide,
                left=phone_left + Inches(0.1), top=phone_top + Inches(2.2),
                width=phone_w - Inches(0.2), height=Inches(1.4),
                text="📱\nวางภาพ screenshot\nที่นี่",
                font_size=14, color=MUTED, align=PP_ALIGN.CENTER)

    add_notes(slide,
              "หน้าผลลัพธ์แสดงคะแนนความน่าเชื่อถือในรูปแบบที่เข้าใจง่าย "
              "พร้อมเหตุผลประกอบเป็นภาษาไทยและรายการแหล่งอ้างอิง "
              "ผู้ใช้สามารถดูรายละเอียดเพิ่มเติมได้ผ่าน accordion sections")


def slide_07a_trending(prs):
    slide = blank_slide(prs)
    add_title(slide, "Demo: หน้า Trending", "FR.6 — ค้นหายอดนิยม (Anonymized)")

    add_bullets(slide,
                left=Inches(0.5), top=Inches(1.4),
                width=Inches(7.5), height=Inches(5),
                items=[
                    "🔥 Rank List — 10 อันดับคำค้นหายอดนิยม",
                    "📅 Period Filter — วันนี้ / สัปดาห์นี้ / เดือนนี้",
                    "🔒 Privacy Notice — ข้อความชี้แจงการ hash query",
                    "👆 Clickable — กดรายการเพื่อส่งไปตรวจสอบทันที",
                    "🔍 SHA-256 Hash — ไม่เก็บ raw text ใน DB",
                ],
                font_size=18)

    phone_left = Inches(8.3)
    phone_top  = Inches(1.1)
    phone_w    = Inches(3.5)
    phone_h    = Inches(5.9)

    frame = slide.shapes.add_shape(1, phone_left, phone_top, phone_w, phone_h)
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)
    frame.line.color.rgb = ACCENT
    frame.line.width = Pt(2.5)

    notch = slide.shapes.add_shape(9, phone_left + Inches(1.2), phone_top + Inches(0.1),
                                   Inches(1.1), Inches(0.25))
    notch.fill.solid()
    notch.fill.fore_color.rgb = ACCENT
    notch.line.fill.background()

    add_textbox(slide,
                left=phone_left + Inches(0.1), top=phone_top + Inches(2.2),
                width=phone_w - Inches(0.2), height=Inches(1.4),
                text="📱\nวางภาพ screenshot\nที่นี่",
                font_size=14, color=MUTED, align=PP_ALIGN.CENTER)

    add_notes(slide,
              "หน้า Trending แสดงหัวข้อที่คนสนใจตรวจสอบมากที่สุด "
              "โดยใช้การ hash ด้วย SHA-256 เพื่อปกป้องความเป็นส่วนตัว "
              "ไม่มีการเก็บข้อความต้นฉบับของผู้ใช้ไว้ใน server")


def slide_07b_history(prs):
    slide = blank_slide(prs)
    add_title(slide, "Demo: หน้าประวัติ (History)", "FR.7 — Local History")

    add_bullets(slide,
                left=Inches(0.5), top=Inches(1.4),
                width=Inches(7.5), height=Inches(5),
                items=[
                    "📋 History List — รายการจาก localStorage ของ browser",
                    "🔍 Filter Bar — ค้นหา + กรองตามประเภท + ระดับคะแนน",
                    "📤 Export JSON — ดาวน์โหลดประวัติเป็น JSON file",
                    "📥 Import JSON — นำเข้าประวัติจาก file",
                    "🗑️ Clear All — ลบประวัติทั้งหมด",
                    "📭 Empty State — แสดงเมื่อยังไม่มีประวัติ",
                ],
                font_size=18)

    phone_left = Inches(8.3)
    phone_top  = Inches(1.1)
    phone_w    = Inches(3.5)
    phone_h    = Inches(5.9)

    frame = slide.shapes.add_shape(1, phone_left, phone_top, phone_w, phone_h)
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)
    frame.line.color.rgb = ACCENT
    frame.line.width = Pt(2.5)

    notch = slide.shapes.add_shape(9, phone_left + Inches(1.2), phone_top + Inches(0.1),
                                   Inches(1.1), Inches(0.25))
    notch.fill.solid()
    notch.fill.fore_color.rgb = ACCENT
    notch.line.fill.background()

    add_textbox(slide,
                left=phone_left + Inches(0.1), top=phone_top + Inches(2.2),
                width=phone_w - Inches(0.2), height=Inches(1.4),
                text="📱\nวางภาพ screenshot\nที่นี่",
                font_size=14, color=MUTED, align=PP_ALIGN.CENTER)

    add_notes(slide,
              "ประวัติการค้นหาทั้งหมดถูกเก็บไว้ใน localStorage ของ browser "
              "ผู้ใช้สามารถส่งออกเป็น JSON เพื่อสำรองข้อมูลหรือนำไปวิเคราะห์เพิ่มเติมได้ "
              "ข้อมูลไม่เคยถูกส่งไปยัง server")


def slide_08_architecture(prs):
    slide = blank_slide(prs)
    add_title(slide, "สถาปัตยกรรมระบบ", "Architecture Diagram")

    # Docker Compose box
    dc_box = slide.shapes.add_shape(1, Inches(0.4), Inches(1.35), Inches(9.5), Inches(5.7))
    dc_box.fill.solid()
    dc_box.fill.fore_color.rgb = RGBColor(0x0D, 0x22, 0x3A)
    dc_box.line.color.rgb = ACCENT
    dc_box.line.width = Pt(1.5)

    add_textbox(slide, Inches(0.6), Inches(1.45), Inches(3), Inches(0.4),
                "🐳 Docker Compose", font_size=14, bold=True, color=ACCENT)

    # Browser box
    add_card(slide, Inches(10.2), Inches(2.5), Inches(2.7), Inches(1.2))
    add_textbox(slide, Inches(10.2), Inches(2.6), Inches(2.7), Inches(0.4),
                "🌐 Browser", font_size=16, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(10.2), Inches(3.0), Inches(2.7), Inches(0.4),
                "Next.js Client\nReact + Zustand", font_size=12, color=MUTED, align=PP_ALIGN.CENTER)

    # Arrow browser -> app
    arr = slide.shapes.add_shape(1, Inches(9.7), Inches(3.0), Inches(0.5), Pt(2))
    arr.fill.solid()
    arr.fill.fore_color.rgb = ACCENT
    arr.line.fill.background()

    add_textbox(slide, Inches(9.8), Inches(2.85), Inches(0.8), Inches(0.3),
                "HTTP", font_size=11, color=MUTED, align=PP_ALIGN.CENTER)

    # Next.js App box
    add_card(slide, Inches(1.0), Inches(1.8), Inches(4.8), Inches(4.7))
    add_textbox(slide, Inches(1.2), Inches(1.9), Inches(4.4), Inches(0.5),
                "⚡ Next.js App  (port 3000)", font_size=17, bold=True, color=ACCENT2)

    apis = [
        "POST /api/analyze/text   — ข้อความ & URL",
        "POST /api/analyze/image  — Gemini AI",
        "POST /api/analyze/source — โดเมน risk",
        "GET  /api/trending       — trending",
        "GET  /api/references/:id — รายการอ้างอิง",
    ]
    add_bullets(slide, Inches(1.2), Inches(2.55), Inches(4.4), Inches(3.5),
                apis, font_size=13, indent_char="")

    # Arrow app -> db
    arr2 = slide.shapes.add_shape(1, Inches(5.8), Inches(3.5), Inches(1.2), Pt(2))
    arr2.fill.solid()
    arr2.fill.fore_color.rgb = ACCENT
    arr2.line.fill.background()
    add_textbox(slide, Inches(5.85), Inches(3.2), Inches(1.1), Inches(0.3),
                "Prisma ORM", font_size=11, color=MUTED, align=PP_ALIGN.CENTER)

    # PostgreSQL box
    add_card(slide, Inches(7.0), Inches(1.8), Inches(2.5), Inches(4.7))
    add_textbox(slide, Inches(7.05), Inches(1.9), Inches(2.4), Inches(0.5),
                "🗄️ PostgreSQL 15\n(port 5432)", font_size=14, bold=True, color=ACCENT2,
                align=PP_ALIGN.CENTER)
    tables = ["Analysis", "Reference", "TrustedSource",
              "KnownFakeClaim", "SuspiciousDomain", "TrendingEntry"]
    add_bullets(slide, Inches(7.1), Inches(2.6), Inches(2.3), Inches(3.5),
                tables, font_size=13, indent_char="▪ ")

    # Gemini external
    add_card(slide, Inches(10.2), Inches(4.2), Inches(2.7), Inches(1.5),
             color=RGBColor(0x1A, 0x35, 0x28))
    add_textbox(slide, Inches(10.2), Inches(4.3), Inches(2.7), Inches(0.4),
                "🤖 Google Gemini API", font_size=14, bold=True, color=GREEN,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(10.2), Inches(4.75), Inches(2.7), Inches(0.5),
                "gemini-2.5-flash\n(image only)", font_size=12, color=MUTED,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(10.2), Inches(5.8), Inches(2.7), Inches(0.3),
                "External API", font_size=11, color=MUTED, align=PP_ALIGN.CENTER)

    add_notes(slide,
              "สถาปัตยกรรมใช้ Docker Compose รวม Next.js App และ PostgreSQL ไว้ด้วยกัน "
              "Gemini API เป็น external dependency เดียวที่ใช้สำหรับวิเคราะห์รูปภาพ "
              "ทุก API endpoint อยู่ใน Next.js App Router")


def slide_09_api(prs):
    slide = blank_slide(prs)
    add_title(slide, "API Design", "5 Endpoints · Zod Validation · Rate Limiting")

    endpoints = [
        ("POST", "/api/analyze/text",    "วิเคราะห์ข้อความและ URL · scoring logic · reference list"),
        ("POST", "/api/analyze/image",   "ส่งรูปไป Gemini · EXIF · AI probability"),
        ("POST", "/api/analyze/source",  "วิเคราะห์โดเมน: อายุ, SSL, WHOIS, phishing"),
        ("GET",  "/api/trending",        "ดึง trending queries (hashed, anonymized)"),
        ("GET",  "/api/references/:id",  "รายการแหล่งอ้างอิงพร้อม stance badge"),
    ]

    for i, (method, path, desc) in enumerate(endpoints):
        cy = Inches(1.45 + i * 0.95)
        add_card(slide, Inches(0.4), cy, Inches(12.3), Inches(0.82))

        method_color = ACCENT if method == "GET" else GREEN
        add_textbox(slide, Inches(0.6), cy + Inches(0.07),
                    Inches(0.8), Inches(0.55),
                    method, font_size=14, bold=True, color=method_color)
        add_textbox(slide, Inches(1.55), cy + Inches(0.07),
                    Inches(3.5), Inches(0.55),
                    path, font_size=15, bold=True, color=TEXT)
        add_textbox(slide, Inches(5.2), cy + Inches(0.07),
                    Inches(7.2), Inches(0.55),
                    desc, font_size=14, color=MUTED)

    # Error format & rate limit
    add_card(slide, Inches(0.4), Inches(6.25), Inches(7.5), Inches(0.9),
             color=RGBColor(0x1A, 0x22, 0x35))
    add_textbox(slide, Inches(0.6), Inches(6.32), Inches(7.2), Inches(0.7),
                '🚫 Error: { error: { code, message_th, message_en } }  —  never raw stack traces',
                font_size=14, color=RED)

    add_card(slide, Inches(8.1), Inches(6.25), Inches(4.6), Inches(0.9),
             color=RGBColor(0x1A, 0x22, 0x35))
    add_textbox(slide, Inches(8.3), Inches(6.32), Inches(4.2), Inches(0.7),
                "⏱️ Rate Limit: 30 req/min per IP (sliding window)",
                font_size=14, color=ORANGE)

    add_notes(slide,
              "ทุก API endpoint ผ่านการ validate ด้วย Zod schema ทั้งฝั่ง client และ server "
              "รูปแบบ error response เป็นมาตรฐานเดียวกันทั้งหมด ไม่มีการ expose stack trace "
              "Rate limiting ใช้ in-memory sliding window 30 req/min ต่อ IP")


def slide_10_scoring(prs):
    slide = blank_slide(prs)
    add_title(slide, "Scoring Logic — คะแนนความน่าเชื่อถือ",
              "Credibility Score Algorithm (0–100)")

    # Score bar
    segments = [
        ("อันตราย\n0–20",   RED,    2.5),
        ("น่าสงสัย\n21–40", ORANGE, 2.5),
        ("ไม่แน่ใจ\n41–60", YELLOW, 2.5),
        ("ค่อนข้างจริง\n61–80", LIME, 2.5),
        ("ยืนยันแล้ว\n81–100", GREEN, 2.5),
    ]
    seg_w = Inches(2.38)
    for i, (label, color, _) in enumerate(segments):
        sx = Inches(0.4 + i * 2.5)
        bar = slide.shapes.add_shape(1, sx, Inches(1.4), seg_w, Inches(0.45))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        add_textbox(slide, sx, Inches(1.9), seg_w, Inches(0.55),
                    label, font_size=12, color=TEXT, align=PP_ALIGN.CENTER)

    # Algorithm steps
    steps = [
        ("🎯", "เริ่มต้น",            "Base score = 50 (neutral)",                   TEXT),
        ("🚨", "Fake claim match",    "พบข้อมูลเท็จที่รู้จัก → score = 5",           RED),
        ("✅", "Trusted source สนับสนุน", "+8 ต่อแหล่ง (สูงสุด 100)",              GREEN),
        ("❌", "Trusted source ต่อต้าน",  "−8 ต่อแหล่ง (ต่ำสุด 0)",               ORANGE),
        ("⚠️", "Suspicious domain",   "−25 (โดเมนน่าสงสัย)",                        ORANGE),
        ("🔍", "Partial match >80%",  "−20 (คล้ายข้อมูลเท็จ > 80%)",              YELLOW),
    ]

    for i, (icon, label, detail, color) in enumerate(steps):
        col = i % 2
        row = i // 2
        cx = Inches(0.4 + col * 6.5)
        cy = Inches(2.65 + row * 1.45)
        add_card(slide, cx, cy, Inches(6.1), Inches(1.25))
        add_textbox(slide, cx + Inches(0.15), cy + Inches(0.1),
                    Inches(0.55), Inches(0.55), icon, font_size=22)
        add_textbox(slide, cx + Inches(0.8), cy + Inches(0.08),
                    Inches(5.0), Inches(0.4),
                    label, font_size=16, bold=True, color=color)
        add_textbox(slide, cx + Inches(0.8), cy + Inches(0.5),
                    Inches(5.0), Inches(0.55),
                    detail, font_size=14, color=MUTED)

    add_notes(slide,
              "อัลกอริทึมการให้คะแนนเริ่มที่ 50 (neutral) แล้วปรับขึ้นหรือลงตามหลักฐาน "
              "การพบ Fake claim ที่รู้จักจะลดคะแนนลงมาก ส่วนแหล่งข่าวน่าเชื่อถือจะเพิ่มหรือลดตาม stance "
              "คะแนนสุดท้ายจะถูก map ไปยัง 5 ระดับ verdict")


def slide_11_gemini(prs):
    slide = blank_slide(prs)
    add_title(slide, "AI Image Analysis — Gemini Integration",
              "FR.5 · gemini-2.5-flash · Graceful Fallback")

    flow = [
        ("1", "Upload", "ผู้ใช้อัปโหลดรูปภาพ"),
        ("2", "Validate", "ตรวจ size/type (Zod)"),
        ("3", "EXIF", "ดึง metadata จากรูป"),
        ("4", "Gemini API", "ส่งไป gemini-2.5-flash\nพร้อม JSON prompt"),
        ("5", "Parse", "แปลง response เป็น\nstructured JSON"),
        ("6", "Response", "aiProbability · detectedModel\nThai reasoning bullets"),
    ]

    arrow_w = Inches(0.35)
    box_w   = Inches(1.8)

    for i, (num, label, detail) in enumerate(flow):
        bx = Inches(0.3 + i * 2.2)
        add_card(slide, bx, Inches(1.5), box_w, Inches(2.2))
        add_textbox(slide, bx, Inches(1.55), box_w, Inches(0.45),
                    num, font_size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_textbox(slide, bx, Inches(2.1), box_w, Inches(0.45),
                    label, font_size=15, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
        add_textbox(slide, bx, Inches(2.65), box_w, Inches(0.85),
                    detail, font_size=12, color=MUTED, align=PP_ALIGN.CENTER)

        if i < len(flow) - 1:
            add_textbox(slide, bx + box_w + Inches(0.05), Inches(2.3),
                        arrow_w, Inches(0.4), "→", font_size=20, color=ACCENT)

    # Fallback card
    add_card(slide, Inches(0.4), Inches(4.0), Inches(12.3), Inches(1.1),
             color=RGBColor(0x1A, 0x28, 0x1A))
    add_textbox(slide, Inches(0.6), Inches(4.08), Inches(12.0), Inches(0.35),
                "🛡️ Graceful Fallback", font_size=16, bold=True, color=GREEN)
    add_textbox(slide, Inches(0.6), Inches(4.5), Inches(12.0), Inches(0.45),
                "หาก Gemini API ล้มเหลว → endpoint คืน 200 พร้อม fallback data "
                "(aiProbability: null, detectedModel: 'unknown', reasoning: ['ไม่สามารถวิเคราะห์ได้'])",
                font_size=14, color=MUTED)

    # Output fields
    add_card(slide, Inches(0.4), Inches(5.3), Inches(12.3), Inches(1.8))
    add_textbox(slide, Inches(0.6), Inches(5.38), Inches(12.0), Inches(0.4),
                "Output Fields:", font_size=15, bold=True, color=ACCENT2)
    out_items = [
        "aiProbability: number (0–1)  — ความน่าจะเป็นที่ถูกสร้างด้วย AI",
        "detectedModel: string        — ชื่อโมเดล AI ที่ตรวจพบ (ถ้ามี)",
        "reasoning: string[]          — เหตุผลประกอบเป็น bullet ภาษาไทย",
    ]
    add_bullets(slide, Inches(0.6), Inches(5.82), Inches(12.0), Inches(1.1),
                out_items, font_size=14, indent_char="• ")

    add_notes(slide,
              "การวิเคราะห์รูปภาพใช้ Gemini 2.5 Flash ซึ่งสามารถตรวจจับรูปภาพที่สร้างด้วย AI "
              "ระบบมีการดึง EXIF metadata ก่อนส่งไป Gemini เพื่อให้ข้อมูลเพิ่มเติม "
              "มีการออกแบบ graceful fallback เพื่อให้ระบบทำงานได้แม้ Gemini API ล้มเหลว")


def slide_12_privacy(prs):
    slide = blank_slide(prs)
    add_title(slide, "Privacy by Design", "ความเป็นส่วนตัวตั้งแต่การออกแบบ")

    decisions = [
        ("💾", "localStorage Only",
         "ประวัติการค้นหาไม่เคยออกจาก browser\nไม่มีการส่งข้อมูลไป server"),
        ("🔐", "SHA-256 Query Hashing",
         "Trending เก็บแค่ hash ของ query\nไม่มี raw text ใน database"),
        ("⏱️", "Rate Limit ไม่ persist",
         "IP ใช้สำหรับ rate limit เท่านั้น\nไม่เก็บ IP ลง database"),
        ("🚫", "Zero Cookies / Sessions",
         "ไม่มี cookie, session, หรือ\nข้อมูลส่วนตัวบน server"),
    ]

    for i, (icon, title_text, desc) in enumerate(decisions):
        col = i % 2
        row = i // 2
        cx = Inches(0.5 + col * 6.4)
        cy = Inches(1.45 + row * 2.3)
        add_card(slide, cx, cy, Inches(6.0), Inches(2.0))
        add_textbox(slide, cx + Inches(0.15), cy + Inches(0.15),
                    Inches(0.7), Inches(0.7), icon, font_size=30)
        add_textbox(slide, cx + Inches(0.95), cy + Inches(0.12),
                    Inches(4.8), Inches(0.5),
                    title_text, font_size=19, bold=True, color=ACCENT)
        add_textbox(slide, cx + Inches(0.95), cy + Inches(0.68),
                    Inches(4.8), Inches(1.0),
                    desc, font_size=15, color=MUTED)

    # Privacy banner
    add_card(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.85),
             color=RGBColor(0x0D, 0x22, 0x1A))
    add_textbox(slide, Inches(0.7), Inches(6.3), Inches(12.0), Inches(0.55),
                "🔒 ข้อมูลเก็บในเครื่องคุณเท่านั้น — ไม่ส่งไป server",
                font_size=18, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    add_notes(slide,
              "Privacy by Design ไม่ใช่แค่ feature แต่เป็นหลักการออกแบบหลักของโปรเจกต์นี้ "
              "ทุก decision ถูกทำขึ้นเพื่อให้มั่นใจว่าข้อมูลส่วนตัวของผู้ใช้ปลอดภัย "
              "โดยไม่กระทบต่อประสิทธิภาพในการวิเคราะห์")


def slide_13_stack(prs):
    slide = blank_slide(prs)
    add_title(slide, "Tech Stack Summary", "เทคโนโลยีที่ใช้และเหตุผล")

    rows = [
        ("Layer",         "Technology",               "Why"),
        ("Framework",     "Next.js 14 (App Router)",  "SSR + API routes ในไฟล์เดียว"),
        ("Language",      "TypeScript (strict)",      "Type safety ลด bugs"),
        ("Styling",       "TailwindCSS + shadcn/ui",  "Rapid UI + accessible components"),
        ("ORM",           "Prisma",                   "Type-safe DB queries + migrations"),
        ("Database",      "PostgreSQL 15",            "Relational + JSONB support"),
        ("State",         "Zustand",                  "Lightweight, no boilerplate"),
        ("Data Fetching", "TanStack Query",            "Caching + loading states"),
        ("Validation",    "Zod",                      "Runtime validation client+server"),
        ("i18n",          "next-intl",                "Thai/English support"),
        ("AI",            "Google Gemini 2.5 Flash",  "Image analysis + Thai reasoning"),
        ("Testing",       "Vitest + Playwright",      "Unit + E2E testing"),
    ]

    col_widths = [Inches(2.0), Inches(3.5), Inches(6.5)]
    col_starts = [Inches(0.4), Inches(2.5), Inches(6.1)]

    for r, row in enumerate(rows):
        cy = Inches(1.35 + r * 0.52)
        is_header = r == 0
        bg = ACCENT if is_header else (SURFACE if r % 2 == 0 else BG)

        for c, (cell, start, width) in enumerate(zip(row, col_starts, col_widths)):
            cell_shape = slide.shapes.add_shape(1, start, cy, width, Inches(0.48))
            cell_shape.fill.solid()
            cell_shape.fill.fore_color.rgb = bg if is_header else (
                RGBColor(0x17, 0x29, 0x40) if r % 2 == 0 else BG)
            cell_shape.line.color.rgb = RGBColor(0x1A, 0x3A, 0x5C)
            cell_shape.line.width = Pt(0.5)

            add_textbox(slide,
                        start + Inches(0.1), cy + Pt(4),
                        width - Inches(0.1), Inches(0.42),
                        cell,
                        font_size=13 if not is_header else 14,
                        bold=is_header,
                        color=BG if is_header else (TEXT if c < 2 else MUTED))

    add_notes(slide,
              "Stack ถูกเลือกโดยให้ความสำคัญกับ Developer Experience และ Type Safety "
              "Next.js App Router ช่วยให้ทำ SSR และ API routes ได้ในโปรเจกต์เดียว "
              "Gemini เป็น external dependency เดียวที่ต้อง API key")


def slide_14_mvp_scope(prs):
    slide = blank_slide(prs)
    add_title(slide, "MVP Scope & Limitations", "จริง vs. Mock ใน MVP")

    col_data = [
        ("✅ จริง (Implemented)", GREEN, [
            "Text & URL analysis",
            "Domain analysis (age, SSL, WHOIS)",
            "Gemini AI image detection",
            "EXIF metadata extraction",
            "Trending (hashed, anonymized)",
            "Local history (localStorage)",
            "i18n TH/EN",
            "Rate limiting (30 req/min)",
        ]),
        ("🟡 Mock (Demo Data)", YELLOW, [
            "Social account analysis",
            "Coordinated behavior detection",
            "Reverse image search",
        ]),
        ("🚫 Out of Scope", RED, [
            "Authentication / user accounts",
            "Production deployment",
            "Real-time updates / WebSocket",
            "Browser extension",
            "Mobile app",
        ]),
    ]

    for i, (header, color, items) in enumerate(col_data):
        cx = Inches(0.4 + i * 4.3)
        add_card(slide, cx, Inches(1.35), Inches(4.1), Inches(5.7))
        add_textbox(slide, cx + Inches(0.15), Inches(1.45), Inches(3.8), Inches(0.55),
                    header, font_size=15, bold=True, color=color)
        add_bullets(slide, cx + Inches(0.15), Inches(2.1), Inches(3.8), Inches(4.8),
                    items, font_size=14, color=TEXT)

    add_notes(slide,
              "MVP มุ่งเน้นฟีเจอร์หลักที่ให้ value จริงแก่ผู้ใช้ "
              "ส่วนที่เป็น mock ถูกออกแบบให้ replace ได้ง่ายด้วย real API ในอนาคต "
              "การไม่มี production deploy เป็นการตัดสินใจตามขอบเขต MVP")


def slide_15_testing(prs):
    slide = blank_slide(prs)
    add_title(slide, "Testing Strategy", "Unit Tests · E2E · Mock Strategy")

    sections = [
        ("🧪 Unit Tests — Vitest", ACCENT, [
            "Services: textAnalyzer, imageAnalyzer, sourceAnalyzer (≥5 tests each)",
            "Hook: useLocalHistory — add, remove, export, import, filter",
            "Zod schemas: validate happy path + edge cases + invalid input",
        ]),
        ("🎭 E2E — Playwright", ACCENT2, [
            "Happy path: submit text → result page with verdict",
            "Happy path: submit URL → result page with references",
            "Trending page: entries load and are clickable",
        ]),
        ("🎭 Mock Strategy — Vitest", MUTED, [
            "vi.mock('../lib/prisma') — mock DB in all service tests",
            "vi.mock('../lib/gemini') — mock Gemini in image tests",
            "No real DB / API calls in unit tests",
        ]),
    ]

    for i, (header, color, items) in enumerate(sections):
        cy = Inches(1.45 + i * 1.9)
        add_card(slide, Inches(0.4), cy, Inches(12.3), Inches(1.7))
        add_textbox(slide, Inches(0.6), cy + Inches(0.1), Inches(12.0), Inches(0.45),
                    header, font_size=17, bold=True, color=color)
        add_bullets(slide, Inches(0.6), cy + Inches(0.65), Inches(12.0), Inches(0.95),
                    items, font_size=14, color=MUTED)

    add_notes(slide,
              "Unit tests ครอบคลุม business logic หลักทั้งหมด "
              "E2E tests ตรวจสอบ happy path ตั้งแต่ต้นจนจบ "
              "การใช้ mock สำหรับ Prisma และ Gemini ทำให้ tests เร็วและ deterministic")


def slide_16_setup(prs):
    slide = blank_slide(prs)
    add_title(slide, "Setup & Running", "3 คำสั่งในการเริ่มใช้งาน")

    cmds = [
        ("1", "docker compose up",
         "เริ่ม Next.js app (port 3000) + PostgreSQL (port 5432)"),
        ("2", "npm run db:seed",
         "Seed ข้อมูล: trusted sources, fake claims, suspicious domains"),
        ("3", "npm test",
         "รัน Vitest unit tests ทั้งหมด"),
    ]

    for i, (num, cmd, desc) in enumerate(cmds):
        cy = Inches(1.5 + i * 1.7)
        add_card(slide, Inches(0.4), cy, Inches(12.3), Inches(1.45))

        # Number badge
        badge = slide.shapes.add_shape(9, Inches(0.55), cy + Inches(0.35),
                                       Inches(0.65), Inches(0.65))
        badge.fill.solid()
        badge.fill.fore_color.rgb = ACCENT
        badge.line.fill.background()
        add_textbox(slide, Inches(0.55), cy + Inches(0.35), Inches(0.65), Inches(0.65),
                    num, font_size=18, bold=True, color=BG, align=PP_ALIGN.CENTER)

        add_textbox(slide, Inches(1.4), cy + Inches(0.1), Inches(10.8), Inches(0.55),
                    cmd, font_size=20, bold=True, color=ACCENT2)
        add_textbox(slide, Inches(1.4), cy + Inches(0.75), Inches(10.8), Inches(0.45),
                    desc, font_size=15, color=MUTED)

    # Env vars
    add_card(slide, Inches(0.4), Inches(6.7), Inches(12.3), Inches(0.6),
             color=RGBColor(0x1A, 0x22, 0x35))
    add_textbox(slide, Inches(0.6), Inches(6.77), Inches(12.0), Inches(0.4),
                "📄 .env.example → ต้องการ: DATABASE_URL  +  GEMINI_API_KEY",
                font_size=15, color=YELLOW)

    add_notes(slide,
              "การ setup ใช้เวลาไม่เกิน 2 นาที เพียงมี Docker และ Node.js "
              "ค่า env ที่จำเป็นมีแค่ DATABASE_URL และ GEMINI_API_KEY "
              "หลัง seed แล้วสามารถเริ่มทดสอบได้ทันที")


def slide_17_roadmap(prs):
    slide = blank_slide(prs)
    add_title(slide, "Future Roadmap (Post-MVP)", "ทิศทางในอนาคต")

    items = [
        ("🔗", "Social Account Analysis API",
         "เชื่อมต่อ API จริงสำหรับตรวจสอบบัญชี social media"),
        ("🖼️", "Reverse Image Search API",
         "ใช้ Google Vision หรือ TinEye API แทน mock"),
        ("🧩", "Browser Extension",
         "ตรวจสอบข่าวได้โดยตรงบน Facebook, LINE Web, Twitter"),
        ("☁️", "Public Deployment + Persistent Accounts",
         "Deploy บน cloud + optional user accounts สำหรับ sync ประวัติ"),
        ("🤝", "Community Fact-Check Database",
         "ระบบ crowdsourcing ให้ผู้ใช้ contribute ข้อมูลเท็จที่พบ"),
    ]

    for i, (icon, title_text, desc) in enumerate(items):
        cy = Inches(1.5 + i * 1.1)
        add_card(slide, Inches(0.4), cy, Inches(12.3), Inches(0.95))
        add_textbox(slide, Inches(0.55), cy + Inches(0.12),
                    Inches(0.6), Inches(0.6), icon, font_size=22)
        add_textbox(slide, Inches(1.25), cy + Inches(0.08),
                    Inches(4.2), Inches(0.45),
                    title_text, font_size=16, bold=True, color=ACCENT)
        add_textbox(slide, Inches(5.6), cy + Inches(0.1),
                    Inches(6.9), Inches(0.65),
                    desc, font_size=14, color=MUTED)

    add_notes(slide,
              "roadmap ถูกจัดลำดับตามผลกระทบที่มีต่อผู้ใช้ "
              "Social API และ Reverse Image Search เป็น quick wins ที่ทำได้ทันที "
              "Community database จะยกระดับความแม่นยำอย่างมากในระยะยาว")


def slide_18_closing(prs):
    slide = blank_slide(prs)

    # Large circle decoration
    circle = slide.shapes.add_shape(9, Inches(8.5), Inches(-1.5), Inches(6), Inches(6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT
    circle.line.fill.background()

    add_textbox(slide,
                left=Inches(0.7), top=Inches(1.5),
                width=Inches(9), height=Inches(1.2),
                text="ชัวร์ก่อนแชร์", font_size=50, bold=True,
                color=ACCENT, align=PP_ALIGN.LEFT)

    add_textbox(slide,
                left=Inches(0.7), top=Inches(2.9),
                width=Inches(9), height=Inches(0.65),
                text="ตรวจสอบข่าวก่อนส่งต่อ — เพราะข่าวปลอมแพร่เร็วกว่าที่คุณคิด",
                font_size=19, color=ACCENT2, align=PP_ALIGN.LEFT)

    line = slide.shapes.add_shape(1, Inches(0.7), Inches(3.75), Inches(4), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    add_textbox(slide,
                left=Inches(0.7), top=Inches(3.95),
                width=Inches(8), height=Inches(0.55),
                text="Thank you · ขอบคุณ", font_size=26, bold=True,
                color=TEXT, align=PP_ALIGN.LEFT)

    add_textbox(slide,
                left=Inches(0.7), top=Inches(4.65),
                width=Inches(8), height=Inches(1.2),
                text="📁 Repository: [github.com/your-org/fake-news]\n"
                     "📧 Contact:    [your.email@example.com]\n"
                     "🐳 Run:        docker compose up",
                font_size=16, color=MUTED, align=PP_ALIGN.LEFT)

    add_textbox(slide,
                left=Inches(0.7), top=Inches(6.1),
                width=Inches(9), height=Inches(0.5),
                text="Q&A — มีคำถามอะไรบ้างครับ/ค่ะ?",
                font_size=20, bold=True, color=ACCENT2, align=PP_ALIGN.LEFT)

    add_notes(slide,
              "ขอบคุณทุกท่านที่รับฟังการนำเสนอ "
              "ชัวร์ก่อนแชร์เป็น MVP ที่พร้อมขยายต่อในอนาคต "
              "ยินดีตอบทุกคำถามเกี่ยวกับสถาปัตยกรรม การออกแบบ หรือการ implement")


# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    slide_01_cover(prs)
    slide_02_problem(prs)
    slide_03_solution(prs)
    slide_04_features(prs)
    slide_05_demo_home(prs)
    slide_06_demo_result(prs)
    slide_07a_trending(prs)
    slide_07b_history(prs)
    slide_08_architecture(prs)
    slide_09_api(prs)
    slide_10_scoring(prs)
    slide_11_gemini(prs)
    slide_12_privacy(prs)
    slide_13_stack(prs)
    slide_14_mvp_scope(prs)
    slide_15_testing(prs)
    slide_16_setup(prs)
    slide_17_roadmap(prs)
    slide_18_closing(prs)

    out_path = os.path.join(os.path.dirname(__file__), "presentation.pptx")
    prs.save(out_path)
    print(f"✅ Saved: {out_path}")
    print(f"   {len(prs.slides)} slides generated")


if __name__ == "__main__":
    main()
